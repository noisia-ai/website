#!/usr/bin/env python3
"""
build-pptx.py — turn a deck spec into an EDITABLE 16:9 PowerPoint/Keynote.

Usage:  python3 build-pptx.py deck.json out.pptx
Requires:  pip install python-pptx

This is the editable-output path (clients who want to tweak slides). The HTML
deck (build-pdf.mjs) stays the source of truth for pixel-perfect PDF; this maps
the same content into native PPTX shapes using the Noisia palette.

deck.json shape:
{
  "title": "Noisia · <Cliente>",
  "slides": [
    {"type": "cover",     "eyebrow": "...", "title": "...", "subtitle": "...", "meta_for": "...", "meta_date": "..."},
    {"type": "statement", "eyebrow": "...", "statement": "..."},
    {"type": "finding",   "eyebrow": "...", "claim": "...", "quote": "...", "source": "...", "confidence": "Alta"},
    {"type": "signal-insight", "eyebrow": "...", "metric": "62%", "metric_label": "...", "quote": "...", "source": "..."},
    {"type": "legal",     "title": "...", "points": ["...", "..."]},
    {"type": "closing",   "statement": "...", "contact": "...", "cta": "..."}
  ]
}
"""
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx")

# Noisia palette (mirror of engine/noisia-tokens.css)
INK = RGBColor(0x0A, 0x0A, 0x0A)
CANVAS = RGBColor(0xFF, 0xFF, 0xFF)
SIGNAL = RGBColor(0x00, 0xEE, 0xEE)
SIGNAL_DARK = RGBColor(0x00, 0x8A, 0x8A)
TENSION = RGBColor(0xEE, 0x0B, 0x00)
BLOB_DARK = RGBColor(0x06, 0x12, 0x18)
EMW, EMH = Inches(13.333), Inches(7.5)  # 16:9


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _text(slide, text, left, top, width, height, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def render(spec, out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMW, EMH
    blank = prs.slide_layouts[6]
    for s in spec.get("slides", []):
        t = s.get("type", "statement")
        slide = prs.slides.add_slide(blank)
        dark = t in ("statement", "closing")
        _bg(slide, BLOB_DARK if dark else CANVAS)
        fg = CANVAS if dark else INK
        if s.get("eyebrow"):
            _text(slide, "· " + s["eyebrow"], Inches(0.9), Inches(0.7), Inches(8), Inches(0.5), 16, SIGNAL_DARK, bold=True)
        if t == "cover":
            _text(slide, s.get("title", ""), Inches(0.9), Inches(2.2), Inches(9), Inches(2.5), 54, fg, bold=True)
            _text(slide, s.get("subtitle", ""), Inches(0.9), Inches(4.6), Inches(8), Inches(1.5), 18, fg)
            _text(slide, f"para · {s.get('meta_for','')}    fecha · {s.get('meta_date','')}",
                  Inches(0.9), Inches(6.6), Inches(11), Inches(0.5), 14, SIGNAL_DARK)
        elif t in ("statement", "closing"):
            _text(slide, s.get("statement", ""), Inches(0.9), Inches(2.6), Inches(11.5), Inches(3), 40, fg, bold=True)
            if s.get("contact"):
                _text(slide, s["contact"], Inches(0.9), Inches(5.8), Inches(11), Inches(0.6), 18, SIGNAL)
        elif t == "finding":
            _text(slide, s.get("claim", ""), Inches(0.9), Inches(1.6), Inches(11.5), Inches(2), 32, fg, bold=True)
            _text(slide, '"' + s.get("quote", "") + '"', Inches(0.9), Inches(3.8), Inches(8), Inches(2), 22, SIGNAL_DARK)
            _text(slide, "— " + s.get("source", ""), Inches(0.9), Inches(5.6), Inches(8), Inches(0.5), 14, fg)
            _text(slide, "confianza: " + s.get("confidence", ""), Inches(9.5), Inches(3.8), Inches(3), Inches(0.6), 20, fg, bold=True)
        elif t == "signal-insight":
            _text(slide, s.get("metric", ""), Inches(0.9), Inches(2.2), Inches(4), Inches(2), 80, INK if not dark else CANVAS, bold=True)
            _text(slide, s.get("metric_label", ""), Inches(0.9), Inches(4.6), Inches(4), Inches(1), 16, SIGNAL_DARK)
            _text(slide, '"' + s.get("quote", "") + '"', Inches(5.2), Inches(2.6), Inches(7), Inches(3), 24, fg)
            _text(slide, "— " + s.get("source", ""), Inches(5.2), Inches(5.6), Inches(7), Inches(0.5), 14, SIGNAL_DARK)
        elif t == "legal":
            _text(slide, s.get("title", "Alcance y confidencialidad"), Inches(0.9), Inches(1.4), Inches(11), Inches(1), 32, fg, bold=True)
            _text(slide, "\n".join("•  " + p for p in s.get("points", [])),
                  Inches(0.9), Inches(2.8), Inches(11.5), Inches(4), 18, fg)
        # footer
        _text(slide, "noisia · social intelligence architects", Inches(0.9), Inches(7.0), Inches(8), Inches(0.4), 11,
              CANVAS if dark else SIGNAL_DARK)
    prs.save(out)
    print(f"Saved {out} ({len(spec.get('slides', []))} slides)")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        sys.exit("Usage: python3 build-pptx.py deck.json out.pptx")
    with open(sys.argv[1], encoding="utf-8") as f:
        render(json.load(f), sys.argv[2])
